from strands import Agent, tool
from strands_tools import use_llm, memory, http_request, generate_image, file_read
from strands.models import BedrockModel
from strands.multiagent.a2a import A2AServer
from strands_tools.a2a_client import A2AClientToolProvider
from botocore.config import Config
import uvicorn
from fastapi import FastAPI
import logging
import os
import sys
from boto3 import Session as AWSSession
import requests
from requests_aws_sign import AWSV4Sign
import websocket
import base64
import threading
import uuid
from typing import Dict, List, Optional, Any, Union
import copy
from opentelemetry import baggage, context
from strands.agent.conversation_manager import SummarizingConversationManager
from shared.response_model import (
    Source,
    SourceSet,
    StructuredDataContent,
    ResponseModel,
)
from shared.image_generator import generate_image_from_descriptions
from shared.adcp_tools import ADCP_TOOLS
import re

# Add the parent directory to path for shared imports
sys.path.append(os.path.join(os.path.dirname(__file__), "..", ".."))
from urllib.parse import urlparse, urlencode, quote, unquote

import argparse
import json
import asyncio
import sys
from bedrock_agentcore.runtime import BedrockAgentCoreApp
from strands.tools.mcp import MCPClient
from shared.short_term_memory_hook import ShortTermMemoryHook
from bedrock_agentcore.memory import MemoryClient

import boto3
from mcp import stdio_client, StdioServerParameters
import re
import base64
import tempfile
from typing import List
from io import BytesIO

# AgentCore Memory Integration
import json
from datetime import datetime

try:
    from bedrock_agentcore.memory import MemoryClient
    from strands.hooks import (
        AgentInitializedEvent,
        HookProvider,
        HookRegistry,
        MessageAddedEvent,
    )

    MEMORY_AVAILABLE = True
except ImportError:
    print("Warning: AgentCore Memory not available, continuing without memory")
    MEMORY_AVAILABLE = False

# AgentCore Memory Integration
from shared.memory_integration import (
    MemoryHookProvider,
    get_memory_configuration,
    create_memory_hooks_and_state,
    extract_session_id_and_memory_id_and_actor_from_payload,
    MEMORY_AVAILABLE,
)

# AgentCore Memory Conversation Manager for persistent session management
from shared.agentcore_memory_conversation_manager import (
    AgentCoreMemoryConversationManager,
    create_agentcore_memory_manager,
)

logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)
logger.addHandler(logging.StreamHandler(sys.stdout))

logging.basicConfig(
    format="%(levelname)s | %(name)s | %(message)s", handlers=[logging.StreamHandler()]
)

os.environ["DEFAULT_TIMEOUT"] = "600"  # set request timeout to 10 minutes
region = os.environ.get("AWS_REGION", "us-east-1")

os.environ["BYPASS_TOOL_CONSENT"] = "true"

knowledgebaseMcpClient = MCPClient(
    lambda: stdio_client(
        StdioServerParameters(
            command="uvx", args=["awslabs.bedrock-kb-retrieval-mcp-server@latest"]
        )
    )
)

# Add the parent directory to path for shared imports
sys.path.append(os.path.join(os.path.dirname(__file__), "..", ".."))

app = BedrockAgentCoreApp()
client = MemoryClient(region_name=os.environ.get("AWS_REGION", "us-east-1"))
memory_name = "Agents_for_Advertising_%s" % datetime.now().strftime("%Y%m%d%H%M%S")

# AppSync Events API configuration
APPSYNC_CHANNEL_NAMESPACE = os.environ.get("APPSYNC_CHANNEL_NAMESPACE", "sessions")
os.environ["AGENT_OBSERVABILITY_ENABLED"] = "true"


def load_external_agents():
    return CONFIG.get("external_agents")


def _normalize_name(name):
    """Normalize filename by lowercasing and replacing separators with a common character."""
    return name.lower().replace("_", "-").replace(" ", "-")


def _find_file_flexible(directory, filename):
    """
    Find a file in directory with flexible matching for underscores, spaces, and hyphens.
    Returns the actual filename if found, None otherwise.
    """
    if not os.path.exists(directory):
        return None

    target_normalized = _normalize_name(filename)

    try:
        files = os.listdir(directory)
        for file in files:
            if _normalize_name(file) == target_normalized:
                return file
    except OSError:
        return None

    return None


def inject_data_into_placeholder(instructions: str, agent_name: str) -> str:
    """
    Replace {{AGENT_NAME}} placeholder with actual agent name in instructions.

    Args:
        instructions: The instruction text containing placeholders
        agent_name: The actual agent name to inject

    Returns:
        Instructions with placeholders replaced
    """
    if not instructions:
        return instructions

    global orchestrator_instance

    # Check if placeholder exists
    if (
        "{{AGENT_NAME}}" not in instructions
        and "{{AGENT_NAME_LIST}}" not in instructions
    ):
        logger.debug(f"No placeholders found in instructions for {agent_name}")
        return instructions

    # Replace {{AGENT_NAME}} placeholder with actual agent name
    if "{{AGENT_NAME}}" in instructions:
        instructions = instructions.replace("{{AGENT_NAME}}", agent_name)
        logger.info(f"✓ Injected agent name '{agent_name}' into instructions")

    # Inject custom values from injectable_values configuration
    injectable_values = get_agent_config(agent_name=agent_name).get(
        "injectable_values", {}
    )
    if injectable_values:
        for key, value in injectable_values.items():
            placeholder = f"{{{{{key}}}}}"  # Creates {{key}} pattern
            if placeholder in instructions:
                instructions = instructions.replace(placeholder, str(value))
                logger.info(
                    f"✓ Injected '{key}' value into instructions for {agent_name}"
                )
        logger.debug(
            f"Processed {len(injectable_values)} injectable values for {agent_name}"
        )

    # Get list of agent names from tool_agent_names and inject into {{AGENT_NAME_LIST}} placeholder
    if "{{AGENT_NAME_LIST}}" in instructions:
        try:
            # Load all agent cards in the agent_cards directory and create a list of objects following the format [{agent_name:string, agent_description:string}], and replace {{AGENT_NAME_LIST}} in the instructions with a strigified version of that list.
            agent_name_list = []
            agent_cards_dir = os.path.join(os.path.dirname(__file__), "agent_cards")
            if os.path.exists(agent_cards_dir):
                for filename in os.listdir(agent_cards_dir):
                    if filename.endswith(".agent.card.json"):
                        card_path = os.path.join(agent_cards_dir, filename)
                        try:
                            with open(card_path, "r") as f:
                                card_data = json.load(f)
                                if (
                                    "agent_name" in card_data
                                    and "agent_description" in card_data
                                ):
                                    agent_name_list.append(
                                        {
                                            "agent_name": card_data["agent_name"],
                                            "agent_description": card_data[
                                                "agent_description"
                                            ],
                                        }
                                    )
                        except Exception as e:
                            logger.warning(f"Failed to load agent card {filename}: {e}")

            if agent_name_list:
                instructions = instructions.replace(
                    "{{AGENT_NAME_LIST}}", json.dumps(agent_name_list)
                )
                logger.info(
                    f"✓ Injected agent name list into instructions: {agent_name_list}"
                )
            else:
                # If no tool agents, replace with empty string or a default message
                instructions = instructions.replace("{{AGENT_NAME_LIST}}", "")
                logger.debug(
                    f"No agent cards found to replace in instructions for {agent_name}. Replaced with empty string"
                )
        except Exception as e:
            logger.error(f"Error injecting agent name list for {agent_name}: {e}")
            # Leave placeholder as-is if there's an error

    return instructions


def load_instructions_for_agent(agent_name: str):
    try:
        base_dir = os.path.dirname(__file__)

        library_dir = os.path.join(base_dir, "agent-instructions-library")

        # Try to list the directory to see what's actually there
        if os.path.exists(library_dir):
            files = os.listdir(library_dir)
        else:
            return "Couldn't load instructions - library directory not found."

        # Try flexible filename matching
        actual_filename = f"{agent_name}.txt"
        instructions_path = os.path.join(library_dir, actual_filename)
        path_exists = os.path.exists(instructions_path)

        if actual_filename and path_exists:
            logging.info(f"Loading instructions from {instructions_path}")
            with open(instructions_path, "r", encoding="utf-8") as f:
                content = f.read().strip()
                # Inject agent name placeholder
                content = inject_data_into_placeholder(content, agent_name)

                return content
        else:
            # File doesn't exist
            logging.warning(f"Instructions file not found: {instructions_path}")
            return "Couldn't load instructions - file not found."

    except FileNotFoundError as e:
        logging.error(f"Warning: instructions.txt not found at {instructions_path}")
        return "Couldn't load instructions."
    except Exception as e:
        logging.error(f"Error loading instructions: {e}")
        return "Couldn't load instructions."


# Load configuration from file
def load_configs(file_name):
    config_path = os.path.join(os.path.dirname(__file__), file_name)
    try:
        with open(config_path, "r", encoding="utf-8") as f:
            return json.load(f)
    except FileNotFoundError:
        print(f"Warning: file not found at {config_path}")
        return {}
    except json.JSONDecodeError as e:
        print(f"Warning: Invalid JSON in {file_name}: {e}")
        return {}


# Load the config once at module level
CONFIG = {}
GLOBAL_CONFIG = {}
response_model_parsed = {}


def get_agent_config(agent_name):
    """Get the configuration for a specific agent"""
    GLOBAL_CONFIG = load_configs("global_configuration.json")
    return GLOBAL_CONFIG.get("agent_configs", {}).get(agent_name, {})


def get_collaborator_agent_model_inputs(agent_name, orchestrator_name):
    """Get the model inputs for a collaborator agent"""
    GLOBAL_CONFIG = load_configs("global_configuration.json")
    orchestrator_config = GLOBAL_CONFIG.get("agent_configs", {}).get(
        orchestrator_name, {}
    )
    model_inputs = orchestrator_config.get("model_inputs", {})
    return model_inputs.get(agent_name, {})


def get_collaborator_agent_config(agent_name, orchestrator_name):
    """Get the configuration for a collaborator agent"""
    GLOBAL_CONFIG = load_configs("global_configuration.json")
    orchestrator_config = GLOBAL_CONFIG.get("agent_configs", {}).get(
        orchestrator_name, {}
    )
    # For collaborators, we need to build a config from the orchestrator's settings
    return {
        "agent_name": agent_name,
        "agent_description": f"Collaborator agent: {agent_name}",
        "model_inputs": orchestrator_config.get("model_inputs", {}).get(agent_name, {}),
    }


# Import shared knowledge base helper
from shared.knowledge_base_helper import (
    setup_agent_knowledge_base,
    get_knowledge_base_tool,
    format_kb_result_with_sources,
    retrieve_knowledge_base_results,
    KnowledgeBaseResult,
    KnowledgeBaseHelper,
)


def load_all_stack_knowledgebase_ids() -> Optional[str]:
    """
    Get knowledge base ID by name pattern (case-insensitive partial match)
    Args:
        name_pattern: Pattern to match against knowledge base names
    Returns:
        str: Knowledge base ID if found, None otherwise
    """
    knowledgebase_ids = {}
    # Discover knowledge bases from AWS
    helper = KnowledgeBaseHelper(logger, os.environ.get("AWS_REGION", "us-east-1"))
    kb_mapping = helper._discover_knowledge_bases_from_aws()
    if not kb_mapping:
        return None
    for kb_name, kb_id in kb_mapping.items():
        if f'{os.environ.get("STACK_PREFIX","XXX")}-' not in kb_name:
            continue
        knowledgebase_ids[kb_name] = kb_id
    return knowledgebase_ids


KNOWLEDGEBASE_IDS = load_all_stack_knowledgebase_ids()


def get_matching_kb_id(name: str) -> Optional[str]:
    return KNOWLEDGEBASE_IDS.get(
        f'{os.environ.get("STACK_PREFIX", "XXX")}-{name}-{os.environ.get("UNIQUE_ID", "XXX")}',
        None,
    )


def get_tool_agent_names():
    """Get the list of tool names that should be wrapped as agent messages"""
    return CONFIG.get("tool_agent_names", [])


all_agents = get_tool_agent_names()
agent_actors = {}
for agent_name in all_agents:
    # Use simple agent name as actor_id to comply with validation pattern
    actor_id = agent_name.replace("_", "-")
    agent_actors[agent_name] = actor_id


from shared.runtime_resolver import RuntimeARNResolver


def get_agentcore_config_from_ssm():
    """
    Retrieve AgentCore configuration from SSM Parameter Store.

    Returns:
        dict: Configuration with agents list containing runtime ARNs and bearer tokens
    """
    try:
        stack_prefix = os.environ.get("STACK_PREFIX", "sim")
        unique_id = os.environ.get("UNIQUE_ID", "")
        region = os.environ.get("AWS_REGION", "us-east-1")

        if not unique_id:
            logging.warning("[get_agentcore_config_from_ssm] UNIQUE_ID not set")
            return None

        parameter_name = f"/{stack_prefix}/agentcore_values/{unique_id}"
        ssm = boto3.client("ssm", region_name=region)

        # Retrieve parameter with decryption
        response = ssm.get_parameter(Name=parameter_name, WithDecryption=True)

        # Parse JSON value
        config_json = response["Parameter"]["Value"]
        config = json.loads(config_json)

        agent_count = len(config.get("agents", []))
        return config

    except Exception as e:
        logging.error(f"[get_agentcore_config_from_ssm] Error: {e}")
        return None


def get_runtime_arn_and_bearer_token(agent_name: str):
    """
    Get runtime URL and bearer token for an agent from SSM Parameter Store or RUNTIMES env var.

    This function retrieves runtime configuration from SSM Parameter Store first,
    then falls back to the RUNTIMES environment variable if SSM is unavailable.

    Args:
        agent_name: Name of the agent to look up

    Returns:
        tuple: (runtime_url, bearer_token) or (None, None) if not found
    """
    region_name = os.environ.get("AWS_REGION", "us-east-1")

    # Try SSM first
    ssm_config = get_agentcore_config_from_ssm()
    if ssm_config:
        agents = ssm_config.get("agents", "{}")
        for agent in agents:
            agent_name_normalized = agent.get("name", "").lower().replace("-", "_")
            search_name_normalized = agent_name.lower().replace("-", "_")
            # Match agent name (handle both hyphen and underscore variations)
            if search_name_normalized in agent_name_normalized:
                runtime_arn = agent.get("runtime_arn", "")
                bearer_token = agent.get("bearer_token", "")

                if runtime_arn:
                    return runtime_arn, bearer_token if bearer_token else ""

    runtimes_env = os.environ.get("RUNTIMES", "")

    if not runtimes_env:
        return None, None

    # Parse RUNTIMES: format is "arn1|token1,arn2|token2,..."
    for entry in runtimes_env.split(","):
        if "|" not in entry:
            # Handle entries without bearer token
            arn = entry
            bearer_token = None
        else:
            arn, bearer_token = entry.split("|", 1)
            bearer_token = bearer_token if bearer_token else None

        # Check if this runtime matches the agent name
        if agent_name.lower().replace("_", "-") in arn.lower():
            print(f"Found runtime for {agent_name} in RUNTIMES env var: {arn[:60]}...")
            if bearer_token:
                print(f"Found bearer token: {bearer_token[:20]}... (truncated)")
            return arn, bearer_token

    print(f"No runtime found for {agent_name}")
    return None, None


# @tool
# async def invoke_external_agent_with_a2a(
#     agent_name: str, prompt: str, session_id: str
# ) -> str:
#     """
#     Invoke an external A2A agent using bearer token authentication.

#     Args:
#         agent_name: Name of the external agent to invoke
#         prompt: The prompt/message to send to the agent

#     Returns:
#         Response from the external agent
#     """
#     print(f"[invoke_external_agent_with_a2a] ===== STARTING A2A INVOCATION =====")
#     print(f"[invoke_external_agent_with_a2a] Invoking external A2A agent: {agent_name}")
#     logging.info(
#         f"[invoke_external_agent_with_a2a] Invoking external A2A agent: {agent_name}"
#     )

#     # Get runtime URL and bearer token from RUNTIMES environment variable
#     print(
#         f"[invoke_external_agent_with_a2a] Calling get_runtime_arn_and_bearer_token..."
#     )
#     runtime_arn, bearer_token = get_runtime_arn_and_bearer_token(agent_name)
#     print(
#         f"[invoke_external_agent_with_a2a] Returned from get_runtime_arn_and_bearer_token"
#     )

#     if not runtime_arn:
#         error_msg = f"Could not find runtime URL for agent: {agent_name}"
#         print(error_msg)
#         return f"Error: {error_msg}"

#     if not bearer_token:
#         error_msg = f"No bearer token found for agent: {agent_name}. Agent may not be A2A-enabled."
#         print(error_msg)
#         return f"Error: {error_msg}"

#     try:
#         # Import the A2A tool creator
#         from shared.a2a_agent_as_tool import send_sync_message

#         # Create A2A tool with bearer token
#         print(f"Creating A2A tool for {agent_name} with bearer token authentication")

#         # Use await instead of asyncio.run() since we're already in an async context
#         response = await send_sync_message(
#             message=prompt,
#             region=os.environ.get("AWS_REGION", "us-east-1"),
#             agent_arn=runtime_arn,
#             bearer_token=bearer_token,
#             session_id=session_id,
#         )
#         print(f"Received response from {agent_name}: {response[:100]}...")
#         return f"<agent-message agent='{agent_name}'>{response}</agent-message>"

#     except Exception as e:
#         error_msg = f"Error invoking A2A agent {agent_name}: {str(e)}"
#         print(error_msg)
#         return f"Error: {error_msg}"


@tool
def invoke_specialist_with_RAG(
    agent_prompt: str, agent_name: str, is_collaborator: bool = True
) -> str:
    global orchestrator_instance
    global collected_sources
    global GLOBAL_CONFIG

    if GLOBAL_CONFIG is None:
        GLOBAL_CONFIG = load_configs("global_configuration.json")
    kb_name = GLOBAL_CONFIG.get("knowledge_bases", {}).get(agent_name)
    if kb_name is not None:
        try:
            logger.info(f"🔧 TOOL: KB name for {agent_name}: {kb_name}")
        except Exception as e:
            print(
                f"agent KB setup failed - this could just be because the agent is not configured to use a knowledgebase: {e}"
            )
    # Get memory configuration from the orchestrator instance if available
    session_id = orchestrator_instance.session_id
    memory_id = orchestrator_instance.memory_id
    orchestrator_name = orchestrator_instance.agent_name
    # Normalize actor_id to comply with validation pattern
    normalized_actor_id = agent_name.replace("_", "-")
    state = {
        "actor_id": normalized_actor_id,
        "session_id": session_id,
        "memory_id": memory_id,
    }
    agent = create_agent(
        agent_name=agent_name, conversation_context="", is_collaborator=True
    )
    result = agent(f"""{agent_prompt}""")

    response_wrapper = (
        f"<agent-message agent='{agent_name}'>{str(result)}</agent-message>"
    )
    return response_wrapper


@tool
def invoke_specialist(agent_prompt: str, agent_name: str) -> str:
    """
    Invoke a specialist agent for collaboration without requiring a knowledge base query.
    
    Use this tool to collaborate with other agents in the agentic advertising ecosystem.
    This is ideal for agents that don't need to query knowledge bases but need to 
    coordinate with other specialists (e.g., VerificationAgent, IdentityAgent, etc.)
    
    Args:
        agent_prompt: The prompt/request to send to the specialist agent
        agent_name: Name of the specialist agent to invoke (e.g., "VerificationAgent", "IdentityAgent")
    
    Returns:
        Response from the specialist agent wrapped in agent-message tags
    """
    global orchestrator_instance
    
    logger.info(f"🔧 TOOL: Invoking specialist agent: {agent_name}")
    
    # Get memory configuration from the orchestrator instance if available
    session_id = orchestrator_instance.session_id
    memory_id = orchestrator_instance.memory_id
    
    # Normalize actor_id to comply with validation pattern
    normalized_actor_id = agent_name.replace("_", "-")
    
    # Create and invoke the specialist agent
    agent = create_agent(
        agent_name=agent_name, conversation_context="", is_collaborator=True
    )
    result = agent(f"""{agent_prompt}""")

    response_wrapper = (
        f"<agent-message agent='{agent_name}'>{str(result)}</agent-message>"
    )
    return response_wrapper


@tool
def lookup_events(agent_name: str, max_results: int = 5) -> str:
    """
    Look up the last things said by a specific agent in the current session.

    Args:
        agent_name: Name of the agent whose events to retrieve
        max_results: Maximum number of events to retrieve (default: 5)

    Returns:
        str: Formatted string containing the agent's recent messages
    """
    global orchestrator_instance

    try:
        # Get session_id and memory_id from orchestrator instance
        session_id = orchestrator_instance.session_id
        memory_id = orchestrator_instance.memory_id

        if not session_id or session_id == "new_session-12345678901234567890":
            return f"No active session found. Cannot retrieve events for {agent_name}."

        # Initialize bedrock-agentcore client
        bedrock_agentcore_client = boto3.client(
            "bedrock-agentcore", region_name=os.environ.get("AWS_REGION", "us-east-1")
        )

        actor_id = agent_name.replace("_", "-")

        logger.info(
            f"🔍 LOOKUP_EVENTS: Looking up events for {agent_name} (actor: {actor_id}) in session {session_id}"
        )

        # Call list_events API
        response = bedrock_agentcore_client.list_events(
            memoryId=memory_id,
            sessionId=session_id,
            actorId=actor_id,
            maxResults=max_results,
        )

        events = response.get("events", [])

        if not events:
            return f"No recent events found for {agent_name} in this session."

        # Try using get_last_k_turns from memory client instead
        # This is more reliable than get_event for retrieving conversation history
        try:
            from bedrock_agentcore.memory import MemoryClient

            memory_client = MemoryClient(
                region_name=os.environ.get("AWS_REGION", "us-east-1")
            )

            recent_turns = memory_client.get_last_k_turns(
                memory_id=memory_id,
                actor_id=actor_id,
                session_id=session_id,
                k=max_results,
                branch_name="main",
                max_results=max_results
                * 2,  # Get more messages to ensure we have enough turns
            )

            if recent_turns:
                result = f"Recent messages from {agent_name}:\n\n"
                for idx, turn in enumerate(recent_turns, 1):
                    for message in turn:
                        role = message.get("role", "unknown")
                        content = message.get("content", {})

                        # Extract text from content
                        if isinstance(content, dict) and "text" in content:
                            text = content["text"]
                        elif isinstance(content, str):
                            text = content
                        else:
                            text = str(content)

                        result += f"{idx}. {role.upper()}:\n{text}\n\n"

                logger.info(
                    f"✅ LOOKUP_EVENTS: Retrieved {len(recent_turns)} turns for {agent_name}"
                )
                return result
            else:
                return f"No recent conversation turns found for {agent_name} in this session."

        except Exception as memory_error:
            logger.warning(
                f"⚠️ LOOKUP_EVENTS: get_last_k_turns failed: {memory_error}, falling back to event-by-event retrieval"
            )

            # Fallback: try get_event for each event
            result = f"Recent messages from {agent_name}:\n\n"

            for idx, event_summary in enumerate(events, 1):
                event_id = event_summary.get("eventId")
                event_type = event_summary.get("eventType", "UNKNOWN")
                timestamp = event_summary.get("timestamp", "N/A")

                # Fetch full event details using get_event
                try:
                    logger.info(
                        f"🔍 LOOKUP_EVENTS: Fetching event {event_id} for {agent_name}"
                    )
                    event_response = bedrock_agentcore_client.get_event(
                        memoryId=memory_id, sessionId=session_id, eventId=event_id
                    )

                    full_event = event_response.get("event", {})
                    logger.info(
                        f"✅ LOOKUP_EVENTS: Got event data: {full_event.keys()}"
                    )

                    # Extract message content based on event type
                    if event_type == "MESSAGE":
                        message_data = full_event.get("message", {})
                        role = message_data.get("role", "unknown")
                        content = message_data.get("content", [])

                        # Extract text from content blocks
                        text_content = []
                        for block in content:
                            if isinstance(block, dict) and "text" in block:
                                text_content.append(block["text"])

                        if text_content:
                            result += f"{idx}. [{timestamp}] {role.upper()}:\n"
                            result += "\n".join(text_content)
                            result += "\n\n"
                    else:
                        # Handle other event types if needed
                        result += f"{idx}. [{timestamp}] Event type: {event_type}\n\n"

                except Exception as event_error:
                    logger.error(
                        f"❌ LOOKUP_EVENTS: Failed to get event {event_id}: {str(event_error)}"
                    )
                    import traceback

                    logger.error(f"   Traceback: {traceback.format_exc()}")
                    result += f"{idx}. [{timestamp}] [Error: {str(event_error)}]\n\n"

            logger.info(
                f"✅ LOOKUP_EVENTS: Retrieved {len(events)} events for {agent_name}"
            )
            return result

    except Exception as e:
        error_msg = f"Error looking up events for {agent_name}: {str(e)}"
        logger.error(f"❌ LOOKUP_EVENTS: {error_msg}")
        return error_msg


@tool
def retrieve_knowledge_base_results_tool(
    agent_name: str, knowledge_base_query: str
) -> str:
    global collected_sources
    global orchestrator_instance
    global response_model_parsed
    kb_name = (
        load_configs("global_configuration.json")
        .get("knowledge_bases", {})
        .get(agent_name)
    )
    logger.info(f"🔧 TOOL: KB name: {kb_name}")
    result_string = "<sources>"
    kb_id = get_matching_kb_id(kb_name)
    if kb_id is None:
        return ""
    else:
        logger.info(f"🔧 TOOL: KB ID: {kb_id}")

    if collected_sources is None:
        collected_sources = {}
    # Pass the full configs dict which contains the knowledge_bases key
    os.environ["STRANDS_KNOWLEDGE_BASE_ID"] = kb_id
    kb_result = retrieve_knowledge_base_results(
        knowledge_base_query,
        agent_name,
        min_score=0.4,
        max_results=3,
        include_metadata=True,
    )

    citations = kb_result.get("citations", [])
    for citation in citations:
        generatedResponse = (
            citation.get("generatedResponsePart", {})
            .get("textResponsePart", {})
            .get("text")
        )
        if "I am unable to" not in generatedResponse:
            result_string += f"<source>{generatedResponse}</source>"
    result_string += "</sources>"
    # generated_text = kb_result.get('output', {}).get('text', '')
    print(f"\n\n\ncitations:\n{result_string}\n\n\n")
    if agent_name not in collected_sources:
        collected_sources[agent_name] = []
    kb_result["query"] = knowledge_base_query
    collected_sources[agent_name].append(kb_result)

    return result_string


def create_agent(agent_name, conversation_context, is_collaborator):
    global orchestrator_instance

    model_inputs = {}
    if is_collaborator:
        model_inputs = get_collaborator_agent_model_inputs(
            agent_name=agent_name, orchestrator_name=orchestrator_instance.agent_name
        )
    else:
        agent_config = get_agent_config(agent_name=agent_name)
        model_inputs = agent_config.get("model_inputs", {}).get(agent_name, {})

    model = BedrockModel(
        model_id=model_inputs.get(
            "model_id", "us.anthropic.claude-sonnet-4-20250514-v1:0"
        ),
        max_tokens=model_inputs.get("max_tokens", 8000),
        top_p=model_inputs.get("top_p", 0.3),
        temperature=model_inputs.get("temperature", 0.8),
        cache_prompt="default",
        cache_tools="default",
    )

    hooks = []
    if "default" in orchestrator_instance.memory_id:
        logger.info(f"🏗️ CREATE_AGENT: Skipping memory hook for default memory_id")
    else:
        # Normalize actor_id to comply with validation pattern
        normalized_actor_id = agent_name.replace("_", "-")
        hooks = [
            ShortTermMemoryHook(
                memory_client=client,
                memory_id=orchestrator_instance.memory_id,
                actor_id=normalized_actor_id,
                session_id=orchestrator_instance.session_id,
            )
        ]

    # Load base instructions and add conversation context if available
    base_instructions = load_instructions_for_agent(agent_name=agent_name)
    enhanced_system_prompt = base_instructions + conversation_context
    from shared.visualization_loader import VisualizationLoader

    loader = VisualizationLoader()
    templates = loader.load_all_templates_for_agent(agent_name)

    # Build visualization context
    if templates:
        viz_context = f"\n\n## Available Visualization Templates for {agent_name}\n\n"
        viz_context += "You have the following visualization templates available:\n\n"

        for template_id, template_data in templates.items():
            usage = template_data.get("usage", "")
            data_mapping = template_data.get("dataMapping", {})

            viz_context += f"### Template: {template_id}\n"
            viz_context += f"**Usage**: {usage}\n\n"
            viz_context += f"**Data Mapping Structure**:\n```json\n{json.dumps(data_mapping, indent=2)}\n```\n\n"

        viz_context += "\n**Instructions**: Map your analysis data to the template fields without modifying the schema, make sure you include the visualizationType and templateId fields in the JSON, and wrap each visualization in the XML tags below:\n"
        viz_context += "<visualization-data type='[template-id]'>[YOUR_MAPPED_JSON_DATA]</visualization-data>\n"
        enhanced_system_prompt += viz_context

    # Build tools list with A2A agent invocation
    tools = [
        # invoke_external_agent_with_a2a,
        retrieve_knowledge_base_results_tool,
        lookup_events,
        file_read,
        generate_image_from_descriptions,
        invoke_specialist_with_RAG,
        http_request,
    ]
    
    # Add AdCP tools for ecosystem agents that need them
    ADCP_ENABLED_AGENTS = [
        "AgencyAgent", "AdvertiserAgent", "PublisherAgent", 
        "SignalAgent", "VerificationAgent", "MeasurementAgent", "IdentityAgent"
    ]
    if agent_name in ADCP_ENABLED_AGENTS:
        tools.extend(ADCP_TOOLS)
        logger.info(f"🔧 CREATE_AGENT: Added AdCP tools for {agent_name}")
    
    collaborator_config = get_collaborator_agent_config(
        agent_name=agent_name, orchestrator_name=orchestrator_instance.agent_name
    )
    if collaborator_config is None:
        collaborator_config = get_agent_config(agent_name=agent_name)

    # Normalize actor_id to comply with validation pattern
    normalized_actor_id = agent_name.replace("_", "-")

    # Create conversation manager for collaborator agents
    # Use AgentCoreMemoryConversationManager when memory is configured
    if "default" not in orchestrator_instance.memory_id.lower():
        conversation_manager = create_agentcore_memory_manager(
            memory_id=orchestrator_instance.memory_id,
            actor_id=normalized_actor_id,
            session_id=orchestrator_instance.session_id,
            region_name=os.environ.get("AWS_REGION", "us-east-1"),
            use_summarizing_fallback=True,
        )
    else:
        conversation_manager = None  # Use default behavior

    return Agent(
        model=model,
        name=agent_name,
        system_prompt=enhanced_system_prompt,
        tools=tools,
        description=collaborator_config.get("agent_description", ""),
        hooks=hooks,
        conversation_manager=conversation_manager,
        state={
            "session_id": orchestrator_instance.session_id,
            "actor_id": normalized_actor_id,
            "memory_id": orchestrator_instance.memory_id,
        },
    )


session = boto3.session.Session()
credentials = session.get_credentials()
appsync_region = session.region_name or os.environ.get("AWS_REGION")

headers = {"Content-Type": "application/json"}
global auth
auth = AWSV4Sign(credentials, appsync_region, "appsync")

# WebSocket connection management
websocket_connections = {}
websocket_lock = threading.Lock()


def transform_response_handler(**event):
    logger.info("transforming response to a structured result")
    yield (ResponseModel.parse_event_loop_structure_to_response_model(event))


def get_websocket_connection(realtime_domain, http_domain):
    """Get or create WebSocket connection for AppSync Events API using IAM authentication"""
    global websocket_connections, websocket_lock

    connection_key = f"{realtime_domain}_{http_domain}"

    with websocket_lock:
        if connection_key in websocket_connections:
            ws = websocket_connections[connection_key]
            if ws and ws.sock and ws.sock.connected:
                return ws
            else:
                # Clean up dead connection
                websocket_connections.pop(connection_key, None)

    try:
        # Use IAM authentication with current session credentials
        global session, credentials
        authorization = {"host": http_domain, "x-amz-user-agent": "aws-sdk-python"}

        # Add AWS credentials to authorization header
        if credentials:
            authorization["authorization"] = (
                f"AWS4-HMAC-SHA256 Credential={credentials.access_key}"
            )
            if credentials.token:
                authorization["x-amz-security-token"] = credentials.token

        header = base64.b64encode(json.dumps(authorization).encode()).decode()
        header = header.replace("+", "-").replace("/", "_").replace("=", "")
        auth_protocol = f"header-{header}"

        # Create WebSocket connection - realtime_domain already includes full domain
        ws_url = f"wss://{realtime_domain}/event/realtime"
        ws = websocket.create_connection(
            ws_url, subprotocols=["aws-appsync-event-ws"], timeout=10
        )

        with websocket_lock:
            websocket_connections[connection_key] = ws

        logger.info(f"✅ CALLBACK: WebSocket connected to {realtime_domain} using IAM")
        return ws

    except Exception as e:
        logger.error(f"❌ CALLBACK: WebSocket connection failed: {e}")
        return None


def set_session_context(session_id):
    """Set the session ID in OpenTelemetry baggage for trace correlation"""
    ctx = baggage.set_baggage("session.id", session_id)
    token = context.attach(ctx)
    logging.info(f"Session ID '{session_id}' attached to telemetry context")
    return token


def appsync_publisher_callback_handler(**kwargs):
    global orchestrator_instance

    # Get AppSync configuration from environment
    appsync_endpoint = os.getenv("APPSYNC_ENDPOINT")
    appsync_realtime_domain = os.getenv("APPSYNC_REALTIME_DOMAIN")
    appsync_channel_namespace = os.getenv("APPSYNC_CHANNEL_NAMESPACE")

    if not appsync_realtime_domain or not appsync_channel_namespace:
        return

    try:
        serializable_data = {}
        for key, value in kwargs.items():
            try:
                json.dumps(value)
                serializable_data[key] = value
            except (TypeError, ValueError):
                if hasattr(value, "__class__"):
                    serializable_data[key] = f"<{value.__class__.__name__}>"
                else:
                    serializable_data[key] = str(value)

        session_id = (
            orchestrator_instance.session_id
            if orchestrator_instance and orchestrator_instance.session_id
            else "default"
        )

        # AppSync Events API WebSocket message
        channel = f"/{appsync_channel_namespace}/{session_id}"
        message = {
            "id": str(uuid.uuid4()),
            "type": "publish",
            "channel": channel,
            "events": [json.dumps(serializable_data)],
        }

        # Extract HTTP domain from endpoint for auth
        http_domain = (
            appsync_endpoint.replace("https://", "").replace("/graphql", "")
            if appsync_endpoint
            else appsync_realtime_domain.replace(
                ".appsync-realtime-api.", ".appsync-api."
            )
        )

        # TODO: publish over http
    except Exception as e:
        logger.error(f"❌ CALLBACK: Failed to publish via WebSocket: {e}")

    # Track tool usage
    if "current_tool_use" in kwargs and kwargs["current_tool_use"].get("name"):
        tool_name = kwargs["current_tool_use"]["name"]
        print(f"🔧 Using tool: {tool_name}")

    if "message" in kwargs and kwargs["message"].get("role") == "assistant":
        print(json.dumps(kwargs["message"], indent=2))


# Global variable to collect sources from tool calls
collected_sources = {}

# Agent context storage for maintaining conversation history across agent switches
# Structure: {session_id: {agent_name: List[messages]}}
# This enables continuous conversation when users switch between agent types
agent_context_store: Dict[str, Dict[str, List[Dict[str, Any]]]] = {}

# Maximum messages to store per agent context (to prevent unbounded memory growth)
MAX_CONTEXT_MESSAGES = 30


def get_context_store_stats() -> Dict[str, Any]:
    """Get statistics about the current context store for debugging."""
    global agent_context_store
    stats = {
        "total_sessions": len(agent_context_store),
        "sessions": {}
    }
    for session_id, agents in agent_context_store.items():
        stats["sessions"][session_id] = {
            "agents": list(agents.keys()),
            "message_counts": {agent: len(msgs) for agent, msgs in agents.items()}
        }
    return stats


def clear_session_context(session_id: str) -> bool:
    """Clear all agent contexts for a specific session."""
    global agent_context_store
    if session_id in agent_context_store:
        del agent_context_store[session_id]
        logger.info(f"🗑️ CONTEXT_CLEAR: Cleared all contexts for session {session_id}")
        return True
    return False


def trim_context_messages(messages: List[Dict[str, Any]], max_messages: int = MAX_CONTEXT_MESSAGES) -> List[Dict[str, Any]]:
    """
    Trim messages to prevent unbounded memory growth.
    Keeps the most recent messages while preserving conversation flow.
    """
    if len(messages) <= max_messages:
        return messages
    
    # Keep the most recent messages
    trimmed = messages[-max_messages:]
    logger.info(f"✂️ CONTEXT_TRIM: Trimmed {len(messages) - max_messages} old messages, keeping {len(trimmed)}")
    return trimmed


@tool
def get_s3_as_base64_and_extract_summary_and_facts(bucket_name, object_key):
    """
    Retrieves a document from S3 and extracts content using appropriate method.

    Args:
        bucket_name (str): The name of the S3 bucket.
        object_key (str): The key (path) of the document in the S3 bucket.

    Returns:
        str: Extracted and analyzed content, or error message.
    """
    s3 = boto3.client("s3")
    try:
        # Get the object from S3
        response = s3.get_object(Bucket=bucket_name, Key=object_key)
        file_content = response["Body"].read()

        # Determine file type from extension
        file_ext = object_key.lower().split(".")[-1]
        logger.info(f"Processing {file_ext} file: {object_key}")

        # Route to appropriate processor
        if file_ext == "pdf":
            return process_pdf_document(file_content, object_key)
        elif file_ext == "docx":
            return process_docx_document(file_content, object_key)
        elif file_ext == "pptx":
            return process_pptx_document(file_content, object_key)
        elif file_ext in ["txt", "md", "csv", "json"]:
            return process_text_document(file_content, object_key)
        elif file_ext in ["png", "jpg", "jpeg", "webp",'gif']:
            return process_image_directly(file_content, object_key,file_ext)
        else:
            return f"Unsupported file type: {file_ext}. Please provide and image or a document of type PDF, DOCX, PPTX, or text files like TXT, MD, CSV, or JSON."

    except Exception as e:
        error_msg = f"Error processing document {object_key}: {str(e)}"
        logger.error(error_msg)
        return error_msg


def process_pdf_document(content: bytes, filename: str) -> str:
    """Process PDF with text extraction and vision fallback"""
    try:
        import PyPDF2
        from io import BytesIO

        # Try text extraction first
        pdf_file = BytesIO(content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)

        extracted_text = ""
        for page in pdf_reader.pages:
            extracted_text += page.extract_text() + "\n"

        # Check if we got meaningful text
        if extracted_text.strip() and len(extracted_text.strip()) > 100:
            logger.info(f"PDF text extraction successful: {len(extracted_text)} chars")
            return analyze_extracted_text(extracted_text, "PDF")
        else:
            # Fall back to vision-based processing
            logger.info("PDF text extraction insufficient, using vision method")
            return process_pdf_with_vision(content, filename)

    except Exception as e:
        logger.error(f"PDF processing error: {e}")
        # Final fallback: base64 encoding
        return process_document_as_base64(content, "PDF")


def process_docx_document(content: bytes, filename: str) -> str:
    """Process DOCX by converting to images and using vision AI"""
    try:
        from docx import Document
        from io import BytesIO

        # Try text extraction first
        doc = Document(BytesIO(content))
        extracted_text = "\n".join(
            [para.text for para in doc.paragraphs if para.text.strip()]
        )

        # Add table content
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                if row_text.strip():
                    extracted_text += f"\n{row_text}"

        if extracted_text.strip() and len(extracted_text.strip()) > 100:
            logger.info(f"DOCX text extraction successful: {len(extracted_text)} chars")
            return analyze_extracted_text(extracted_text, "Word document")
        else:
            # Fall back to vision processing
            logger.info("DOCX text extraction insufficient, using vision method")
            images = convert_docx_to_images(content)
            return process_images_with_vision(images, filename, "docx")

    except Exception as e:
        logger.error(f"DOCX processing error: {e}")
        return f"Error processing Word document: {str(e)}"


def process_pptx_document(content: bytes, filename: str) -> str:
    """Process PPTX by converting to images and using vision AI"""
    try:
        from pptx import Presentation
        from io import BytesIO

        # Try text extraction first
        prs = Presentation(BytesIO(content))
        extracted_text = ""

        for i, slide in enumerate(prs.slides):
            extracted_text += f"\n\n=== Slide {i + 1} ===\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    extracted_text += f"{shape.text}\n"

        if extracted_text.strip() and len(extracted_text.strip()) > 100:
            logger.info(f"PPTX text extraction successful: {len(extracted_text)} chars")
            return analyze_extracted_text(extracted_text, "PowerPoint presentation")
        else:
            # Fall back to vision processing
            logger.info("PPTX text extraction insufficient, using vision method")
            images = convert_pptx_to_images(content)
            return process_images_with_vision(images, filename, "pptx")

    except Exception as e:
        logger.error(f"PPTX processing error: {e}")
        return f"Error processing PowerPoint: {str(e)}"


def process_text_document(content: bytes, filename: str) -> str:
    """Process plain text documents"""
    try:
        text = content.decode("utf-8")
        logger.info(f"Text file processed: {len(text)} chars")
        return analyze_extracted_text(text, "text file")
    except Exception as e:
        logger.error(f"Text processing error: {e}")
        return f"Error processing text file: {str(e)}"
    
def process_image_directly(content: bytes, filename: str, filetype:str) -> str:
    """Process plain text documents"""
    try:
        return process_single_image_with_vision(content,0,filetype)
    except Exception as e:
        logger.error(f"Text processing error: {e}")
        return f"Error processing text file: {str(e)}"


def analyze_extracted_text(text: str, doc_type: str) -> str:
    """Analyze extracted text using Bedrock"""
    try:
        model = BedrockModel(
            model_id="us.anthropic.claude-sonnet-4-20250514-v1:0",
            max_tokens=8000,
            top_p=0.8,
            temperature=0.3,
        )
        analysisAgent = Agent(
            model=model,
            system_prompt=f"You are an expert in analyzing {doc_type} content. Extract key facts, insights, and summarize the main points concisely.",
            tools=[],
            description="Document analysis agent",
        )

        # Limit text length for analysis
        text_to_analyze = text[:15000] if len(text) > 15000 else text
        analysis = analysisAgent(
            f"Analyze this {doc_type} content:\n\n{text_to_analyze}"
        )
        return str(analysis)

    except Exception as e:
        logger.error(f"Text analysis error: {e}")
        return f"Extracted content (analysis failed): {text[:1000]}..."


def process_pdf_with_vision(content: bytes, filename: str) -> str:
    """Process PDF using vision AI for image-based PDFs"""
    try:
        base64_pdf = base64.b64encode(content).decode("utf-8")

        model = BedrockModel(
            model_id="us.anthropic.claude-sonnet-4-20250514-v1:0",
            max_tokens=8000,
            top_p=0.8,
            temperature=0.3,
        )
        analysisAgent = Agent(
            model=model,
            system_prompt="You are an expert in extracting content from document images. Describe all visible text, charts, and visual elements.",
            tools=[],
            description="Vision-based document analyzer",
        )

        analysis = analysisAgent(f"Analyze this PDF document: {base64_pdf}")
        return str(analysis)

    except Exception as e:
        logger.error(f"Vision PDF processing error: {e}")
        return f"Error processing PDF with vision: {str(e)}"


def process_document_as_base64(content: bytes, doc_type: str) -> str:
    """Final fallback: process document as base64"""
    try:
        base64_content = base64.b64encode(content).decode("utf-8")
        return f"Document encoded as base64 ({doc_type}): {base64_content[:200]}... [truncated]"
    except Exception as e:
        return f"Error encoding document: {str(e)}"


def convert_docx_to_images(content: bytes) -> List:
    """Convert DOCX to images for vision processing"""
    try:
        from docx import Document
        from PIL import Image, ImageDraw, ImageFont
        from io import BytesIO
        import tempfile

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as temp_file:
            temp_file.write(content)
            temp_file.flush()  # Ensure content is written to disk
            temp_path = temp_file.name

        try:
            doc = Document(temp_path)
            images = []

            # Create text-based images from content
            current_content = []
            for para in doc.paragraphs:
                if para.text.strip():
                    current_content.append(para.text)
                    if len(current_content) >= 20:  # Create image every 20 lines
                        img = create_text_image(current_content)
                        images.append(img)
                        current_content = []

            if current_content:
                img = create_text_image(current_content)
                images.append(img)

            return images
        finally:
            import os

            if os.path.exists(temp_path):
                os.unlink(temp_path)

    except Exception as e:
        logger.error(f"DOCX to images conversion error: {e}")
        return []


def convert_pptx_to_images(content: bytes) -> List:
    """Convert PPTX to images for vision processing"""
    try:
        from pptx import Presentation
        from io import BytesIO
        import tempfile

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
            temp_file.write(content)
            temp_file.flush()  # Ensure content is written to disk
            temp_path = temp_file.name

        try:
            prs = Presentation(temp_path)
            images = []

            for i, slide in enumerate(prs.slides):
                slide_content = [f"=== Slide {i + 1} ==="]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)

                if slide_content:
                    img = create_text_image(slide_content, slide_format=True)
                    images.append(img)

            return images
        finally:
            import os

            if os.path.exists(temp_path):
                os.unlink(temp_path)

    except Exception as e:
        logger.error(f"PPTX to images conversion error: {e}")
        return []


def create_text_image(text_lines: List[str], slide_format: bool = False):
    """Create PIL Image from text content"""
    try:
        from PIL import Image, ImageDraw, ImageFont

        width = 1200 if slide_format else 800
        height = max(600, len(text_lines) * 30 + 100)

        img = Image.new("RGB", (width, height), color="white")
        draw = ImageDraw.Draw(img)

        try:
            font = ImageFont.load_default()
        except:
            font = None

        y_offset = 20
        line_height = 22

        for line in text_lines:
            if not line:
                y_offset += line_height // 2
                continue

            # Wrap long lines
            if len(line) > 90:
                words = line.split()
                current_line = ""
                for word in words:
                    test_line = current_line + word + " "
                    if len(test_line) < 90:
                        current_line = test_line
                    else:
                        if current_line.strip():
                            if font:
                                draw.text(
                                    (20, y_offset),
                                    current_line.strip(),
                                    font=font,
                                    fill="black",
                                )
                            else:
                                draw.text(
                                    (20, y_offset), current_line.strip(), fill="black"
                                )
                            y_offset += line_height
                        current_line = word + " "

                if current_line.strip():
                    if font:
                        draw.text(
                            (20, y_offset),
                            current_line.strip(),
                            font=font,
                            fill="black",
                        )
                    else:
                        draw.text((20, y_offset), current_line.strip(), fill="black")
                    y_offset += line_height
            else:
                if font:
                    draw.text((20, y_offset), line, font=font, fill="black")
                else:
                    draw.text((20, y_offset), line, fill="black")
                y_offset += line_height

        return img

    except Exception as e:
        logger.error(f"Text image creation error: {e}")
        # Return minimal fallback image
        from PIL import Image

        return Image.new("RGB", (400, 100), color="white")


def process_images_with_vision(images: List, filename: str, doc_type: str) -> str:
    """Process images using Claude vision"""
    try:
        if not images:
            return f"No images generated from {doc_type} file"

        all_text = []
        for i, image in enumerate(images):
            try:
                page_text = process_single_image_with_vision(image, i + 1, doc_type)
                if page_text.strip():
                    all_text.append(f"=== Page/Slide {i + 1} ===\n{page_text}")
            except Exception as e:
                logger.error(f"Error processing image {i + 1}: {e}")
                all_text.append(f"=== Page/Slide {i + 1} ===\n[Error: {str(e)}]")

        combined_text = "\n\n".join(all_text)
        return analyze_extracted_text(combined_text, doc_type)

    except Exception as e:
        logger.error(f"Vision processing error: {e}")
        return f"Error processing with vision: {str(e)}"


def process_single_image_with_vision(image, page_num: int, doc_type: str) -> str:
    """Process single image with Claude vision"""
    try:
        bedrock_runtime = boto3.client("bedrock-runtime")

        prompt = ''
        response = []
        if doc_type in ['png','jpeg','jpg','webp','gif']:
            if doc_type == 'jpg':
                doc_type = 'jpeg'
            prompt = f"Analyze this image. If it appears to be a document, extract ALL visible text and describe any visual elements, charts, or diagrams. If it appears to be visual imagery or a creative asset for an ad campaign, return a detailed description of what the image depicts, also noting sentiment, mood, and artistic techniques."
            response = bedrock_runtime.converse(
               modelId="global.anthropic.claude-sonnet-4-5-20250929-v1:0",
               messages=[
                   {
                       "role": "user",
                       "content": [
                           {"text": prompt},
                           {"image": {"format": doc_type, "source": {"bytes": image}}},
                       ],
                   }
               ],
               inferenceConfig={"maxTokens": 12000, "temperature": 0.8},
           )
        else: 
            from io import BytesIO

            # Convert PIL Image to bytes
            buffer = BytesIO()
            image.save(buffer, format="PNG", optimize=True, quality=85)
            image_bytes = buffer.getvalue()

        
            prompt = f"Analyze this {doc_type} page/slide {page_num}. Extract ALL visible text and describe any visual elements, charts, or diagrams."

            response = bedrock_runtime.converse(
                modelId="global.anthropic.claude-sonnet-4-5-20250929-v1:0",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"text": prompt},
                            {"image": {"format": "png", "source": {"bytes": image_bytes}}},
                        ],
                    }
                ],
                inferenceConfig={"maxTokens": 12000, "temperature": 0.1},
            )

        response_text = ""
        if "output" in response and "message" in response["output"]:
            content = response["output"]["message"]["content"]
            for item in content:
                if "text" in item:
                    response_text += item["text"]

        return (
            response_text.strip()
            if response_text
            else f"[No content from page {page_num}]"
        )

    except Exception as e:
        logger.error(f"Single image vision error: {e}")
        return f"[Vision processing failed: {str(e)[:100]}]"


class GenericAgent:
    def __init__(self):
        self.logger = logger or logging.getLogger(__name__)
        self.region = region or os.environ.get("AWS_DEFAULT_REGION", "us-east-1")

        # Initialize memory-related properties
        self.bedrock_client = boto3.client("bedrock-runtime")

        # Initialize conversation context for memory integration
        self.conversation_context = ""
        self.memory_id = os.environ.get("MEMORY_ID", "")
        self.session_id = "new_session"

        # Initialize sources collection
        global collected_sources
        collected_sources = {}
        global appsync_publisher_callback_handler
        global CONFIG
        # Create the summarizing conversation manager with default settings

    def create_orchestrator(self, session_id, memory_id, agent_name, saved_messages: Optional[List[Dict[str, Any]]] = None):
        """
        Create an orchestrator agent with optional conversation history restoration.
        
        Args:
            session_id: The session identifier (shared across all agents in a conversation)
            memory_id: The AgentCore memory identifier
            agent_name: Name of the agent to create
            saved_messages: Optional list of previous messages to restore conversation context
        """
        if agent_name == "default":
            return Agent()
        self.agent_name = agent_name
        # Load configuration
        try:
            config = get_agent_config(agent_name=agent_name)
            self.team_name = config.get("team_name", "")
        except Exception as e:
            config = {
                "agent_name": agent_name,
                "agent_display_name": agent_name,
                "agent_description": "Default agent description",
                "team_name": "Default team",
                "use_handler_template": True,
                "tool_agent_names": [],
                "external_agents": [],
                "model_inputs": {
                    f"{agent_name}": {
                        "model_id": "us.anthropic.claude-sonnet-4-20250514-v1:0",
                        "max_tokens": 12000,
                        "temperature": 0.3,
                        "top_p": 0.8,
                    }
                },
            }

        # Extract model inputs
        try:
            model_inputs = config.get("model_inputs", {}).get(agent_name, {})
        except Exception as e:
            model_inputs = {
                "model_id": "us.anthropic.claude-sonnet-4-20250514-v1:0",
                "max_tokens": 12000,
                "temperature": 0.3,
                "top_p": 0.8,
            }

        try:
            model = BedrockModel(
                model_id=model_inputs.get(
                    "model_id", "us.anthropic.claude-sonnet-4-20250514-v1:0"
                ),
                max_tokens=model_inputs.get("max_tokens", 12000),
                cache_prompt="default",
                cache_tools="default",
            )
            if model_inputs.get("temperature"):
                model.temperature = model_inputs.get("temperature")
            if model_inputs.get("top_p"):
                model.top_p = model_inputs.get("top_p")
        except Exception as e:
            logger.error(f"✗ Failed to create Bedrock model: {e}")
            import traceback

            logger.error(f"   Traceback: {traceback.format_exc()}")
        # Setup memory hooks
        hooks = []
        self.session_id = session_id
        self.memory_id = memory_id

        if "default" in memory_id:
            logger.info(f"⊘ Skipping memory hook for default memory_id")
        else:
            try:
                # Normalize actor_id to comply with validation pattern
                normalized_actor_id = agent_name.replace("_", "-")
                hooks = [
                    ShortTermMemoryHook(
                        client, memory_id, normalized_actor_id, session_id
                    )
                ]
            except Exception as e:
                logger.error(f"✗ Failed to create memory hook: {e}")
                import traceback

                logger.error(f"   Traceback: {traceback.format_exc()}")

        # Load instructions
        logger.info(f"\n📝 Loading agent instructions...")
        print(f"\n📝 Loading agent instructions...")
        try:
            base_instructions = load_instructions_for_agent(agent_name=agent_name)
            from shared.visualization_loader import VisualizationLoader

            loader = VisualizationLoader()
            templates = loader.load_all_templates_for_agent(agent_name)

            # Build visualization context
            if templates:
                viz_context = (
                    f"\n\n## Available Visualization Templates for {agent_name}\n\n"
                )
                viz_context += (
                    "You have the following visualization templates available:\n\n"
                )

                for template_id, template_data in templates.items():
                    usage = template_data.get("usage", "")
                    data_mapping = template_data.get("dataMapping", {})

                    viz_context += f"### Template: {template_id}\n"
                    viz_context += f"**Usage**: {usage}\n\n"
                    viz_context += f"**Data Mapping Structure**:\n```json\n{json.dumps(data_mapping, indent=2)}\n```\n\n"

                viz_context += "\n**Instructions**: Map your analysis data to the appropriate template fields without modifying the schema and wrap each visualization in the XML tags below:\n"
                viz_context += "<visualization-data type='[template-id]'>[YOUR_MAPPED_JSON_DATA]</visualization-data>\n"
                base_instructions = base_instructions + viz_context
            instruction_length = len(base_instructions) if base_instructions else 0
            if instruction_length == 0:
                logger.warning(f"⚠️  WARNING: Instructions are empty!")
        except Exception as e:
            logger.error(f"✗ Failed to load instructions: {e}")
            import traceback

            logger.error(f"   Traceback: {traceback.format_exc()}")

        try:
            enhanced_system_prompt = base_instructions + self.conversation_context
        except Exception as e:
            logger.error(f"✗ Failed to build system prompt: {e}")
            print(f"✗ Failed to build system prompt: {e}")
        try:
            tools = [
                invoke_specialist_with_RAG,
                retrieve_knowledge_base_results_tool,
                generate_image_from_descriptions,
                lookup_events,
            ]
            
            # Add AdCP tools for ecosystem agents that need them
            ADCP_ENABLED_AGENTS = [
                "AgencyAgent", "AdvertiserAgent", "PublisherAgent", 
                "SignalAgent", "VerificationAgent", "MeasurementAgent", "IdentityAgent"
            ]
            if agent_name in ADCP_ENABLED_AGENTS:
                tools.extend(ADCP_TOOLS)
                logger.info(f"🔧 CREATE_ORCHESTRATOR: Added AdCP tools for {agent_name}")
        except Exception as e:
            logger.error(f"✗ Failed to build tools list: {e}")

        # Check environment variables
        appsync_endpoint = os.getenv("APPSYNC_ENDPOINT")
        appsync_vars = [k for k in os.environ.keys() if "APPSYNC" in k.upper()]

        try:
            agent_description = config.get("agent_description", "")

            # Normalize actor_id to comply with validation pattern
            actor_name = agent_name or config.get("agent_id", config.get("agent_name"))
            normalized_actor_id = actor_name.replace("_", "-")

            # Create AgentCore Memory Conversation Manager for persistent session management
            # This provides:
            # 1. Automatic persistence of conversation history to AgentCore Memory
            # 2. Retrieval of conversation history when resuming sessions (survives process restarts)
            # 3. Falls back to SummarizingConversationManager for context window management
            if "default" not in self.memory_id.lower():
                conversation_manager = create_agentcore_memory_manager(
                    memory_id=self.memory_id,
                    actor_id=normalized_actor_id,
                    session_id=self.session_id,
                    region_name=os.environ.get("AWS_REGION", "us-east-1"),
                    use_summarizing_fallback=True,  # Use SummarizingConversationManager for context reduction
                )
                logger.info(
                    f"✅ Using AgentCoreMemoryConversationManager for {agent_name} "
                    f"(memory_id={self.memory_id}, session_id={self.session_id})"
                )
            else:
                # Fall back to simple SummarizingConversationManager when memory is not configured
                conversation_manager = SummarizingConversationManager(
                    summary_ratio=0.3,
                    preserve_recent_messages=5,
                    summarization_system_prompt="summarize the current conversation context.",
                )
                logger.info(f"⊘ Using SummarizingConversationManager (no memory configured) for {agent_name}")

            # Build agent kwargs
            agent_kwargs = {
                "model": model,
                "name": actor_name,
                "system_prompt": enhanced_system_prompt,
                "tools": tools,
                "description": agent_description,
                "hooks": hooks,
                "state": {
                    "session_id": self.session_id,
                    "actor_id": normalized_actor_id,
                    "memory_id": self.memory_id,
                },
                "conversation_manager": conversation_manager,
            }
            
            # Restore conversation history if provided
            if saved_messages:
                agent_kwargs["messages"] = saved_messages
                logger.info(f"📂 CONTEXT_RESTORE: Restored {len(saved_messages)} messages for {agent_name}")
            
            agent = Agent(**agent_kwargs)
            return agent
        except Exception as e:
            logger.error(f"CREATE_ORCHESTRATOR: FAILED")
            logger.error(f"Error: {e}")
            import traceback

            logger.error(f"\nFull traceback:")
            logger.error(traceback.format_exc())
            logger.error(f"{'='*80}\n")


print(f"DEBUG: Module load - All env vars: {list(os.environ.keys())}")

# Create the orchestrator instance
orchestrator_instance = GenericAgent()
agent = None  # Will be lazily initialized on first invocation
current_agent_name = None  # Track which agent is currently loaded


def _flush_log(message: str, level: str = "INFO"):
    """Force-flush a log message to ensure it appears in CloudWatch immediately."""
    import sys
    try:
        timestamp = datetime.now().isoformat()
        formatted = f"[{timestamp}] {level}: {message}"
        print(formatted, flush=True)
        sys.stdout.flush()
        sys.stderr.flush()
        if level == "ERROR":
            logger.error(message)
        elif level == "WARNING":
            logger.warning(message)
        elif level == "DEBUG":
            logger.debug(message)
        else:
            logger.info(message)
    except Exception as log_err:
        # Last resort - just print
        print(f"LOG_ERROR: {log_err} | Original message: {message}", flush=True)


@app.entrypoint
async def agent_invocation(payload, context):
    """
    Invoke the orchestrator, unless directed otherwise
    Returns complete response chunks instead of streaming tokens
    """
    _flush_log("=" * 60)
    _flush_log("🚀 AGENT_INVOCATION: Entry point called")
    _flush_log(f"🚀 AGENT_INVOCATION: Payload keys: {list(payload.keys()) if payload else 'None'}")
    _flush_log(f"🚀 AGENT_INVOCATION: Context type: {type(context)}")
    
    try:
        appsync_endpoint = os.getenv("APPSYNC_ENDPOINT")
        _flush_log(f"🔧 AGENT_INVOCATION: APPSYNC_ENDPOINT = {appsync_endpoint[:50] if appsync_endpoint else 'None'}...")

        # Check for APPSYNC variations
        for var_name in [
            "APPSYNC_ENDPOINT",
            "AppSyncEndpoint",
            "appsync_endpoint",
            "APPSYNC_URL",
        ]:
            value = os.getenv(var_name)

        global collected_sources
        global agent
        global CONFIG
        global GLOBAL_CONFIG
        global memory_id
        global orchestrator_instance
        global current_agent_name
        
        _flush_log(f"🔧 AGENT_INVOCATION: Global vars accessed. current_agent_name={current_agent_name}, agent={'exists' if agent else 'None'}")

        # Process the prompt - it can be a string or a list of content blocks
        raw_prompt = payload.get("prompt")
        _flush_log(f"📝 AGENT_INVOCATION: raw_prompt type={type(raw_prompt)}, length={len(str(raw_prompt)) if raw_prompt else 0}")
        media = payload.get("media", {})

        # Parse the prompt to extract text and file attachments
        user_input = ""
        file_attachments = []
        seen_files = set()  # Track file names to avoid duplicates
        _flush_log(f"📝 AGENT_INVOCATION: Parsing prompt...")

        if isinstance(raw_prompt, str):
            # Simple string prompt
            user_input = raw_prompt
            _flush_log(f"📝 AGENT_INVOCATION: Prompt is string, length={len(user_input)}")
        elif isinstance(raw_prompt, list):
            _flush_log(f"📝 AGENT_INVOCATION: Prompt is list with {len(raw_prompt)} blocks")
            # Content blocks format (Bedrock format)
            for idx, block in enumerate(raw_prompt):
                _flush_log(f"📝 AGENT_INVOCATION: Processing block {idx}: type={type(block)}")
                if isinstance(block, dict):
                    if "text" in block:
                        # Extract text content
                        text_content = block["text"]
                        user_input = text_content
                        _flush_log(f"📝 AGENT_INVOCATION: Found text block, length={len(text_content)}")
                    if "document" in block:
                        # Extract document content
                        _flush_log(f"📝 AGENT_INVOCATION: Found document block")
                        document_content = block["document"]
                        file_attachments.append(document_content)
                elif isinstance(block, str):
                    user_input += block
        else:
            _flush_log(f"📝 AGENT_INVOCATION: Prompt is other type: {type(raw_prompt)}")
            user_input = str(raw_prompt)
        _flush_log(f"📝 AGENT_INVOCATION: user_input length={len(user_input)}, file_attachments={len(file_attachments)}")
        
        # Process file attachments - download from S3 and convert to bytes
        processed_attachments = []
        s3_client = boto3.client(
            "s3", region_name=os.environ.get("AWS_REGION", "us-east-1")
        )

        # Check for explicit direct mention flag from frontend (clean approach)
        direct_mention_target = payload.get("direct_mention_target")
        direct_mention_mode = False

        # Clear collected sources from previous invocations
        collected_sources = {}

        # Extract session information from payload for memory integration
        _flush_log("🔍 AGENT_INVOCATION: Extracting session info from payload...")
        try:
            session_id, extracted_memory_id, agent_name = (
                extract_session_id_and_memory_id_and_actor_from_payload(payload)
            )
            _flush_log(f"🔍 AGENT_INVOCATION: session_id={session_id}, memory_id={extracted_memory_id}, agent_name={agent_name}")
        except Exception as extract_err:
            _flush_log(f"❌ AGENT_INVOCATION: Failed to extract session info: {extract_err}", "ERROR")
            import traceback
            _flush_log(f"❌ AGENT_INVOCATION: Traceback: {traceback.format_exc()}", "ERROR")
            raise

        _flush_log("📂 AGENT_INVOCATION: Loading global configuration...")
        GLOBAL_CONFIG = load_configs("global_configuration.json")
        _flush_log(f"📂 AGENT_INVOCATION: GLOBAL_CONFIG keys: {list(GLOBAL_CONFIG.keys()) if GLOBAL_CONFIG else 'None'}")

        _flush_log(f"📂 AGENT_INVOCATION: Getting agent config for {agent_name}...")
        CONFIG = get_agent_config(agent_name=agent_name)
        _flush_log(f"📂 AGENT_INVOCATION: CONFIG keys: {list(CONFIG.keys()) if CONFIG else 'None'}")

        # Check if we need to create or recreate the agent
        agent_type_changed = (current_agent_name is not None and current_agent_name != agent_name)
        _flush_log(f"🔄 AGENT_INVOCATION: agent_type_changed={agent_type_changed}, current={current_agent_name}, new={agent_name}")
        
        # Access global context store
        global agent_context_store
        
        # Determine the session key for context storage (use session_id, shared across all agents)
        context_session_key = session_id if session_id else "default_session"
        _flush_log(f"🔄 AGENT_INVOCATION: context_session_key={context_session_key}")
        
        if agent is None or agent_type_changed:
            _flush_log(f"🏗️ AGENT_INVOCATION: Need to create agent (agent is None: {agent is None}, type_changed: {agent_type_changed})")
            # Need to create a new agent (first time or agent type changed)
            
            # SAVE current agent's context before switching (if switching)
            if agent_type_changed and agent is not None and current_agent_name:
                # Initialize session entry if needed
                if context_session_key not in agent_context_store:
                    agent_context_store[context_session_key] = {}
                
                # Deep copy messages to preserve state
                try:
                    saved_messages = copy.deepcopy(agent.messages) if hasattr(agent, 'messages') and agent.messages else []
                    # Trim to prevent unbounded memory growth
                    saved_messages = trim_context_messages(saved_messages)
                    agent_context_store[context_session_key][current_agent_name] = saved_messages
                    _flush_log(f"💾 CONTEXT_SAVE: Saved {len(saved_messages)} messages for {current_agent_name}")
                except Exception as e:
                    _flush_log(f"⚠️ CONTEXT_SAVE: Failed to save context for {current_agent_name}: {e}", "WARNING")
            
            if agent_type_changed:
                _flush_log(f"🔄 Agent type changed from {current_agent_name} to {agent_name}")
            
            # Set up session info
            if session_id:
                orchestrator_instance.session_id = session_id
                orchestrator_instance.memory_id = extracted_memory_id
                orchestrator_instance.direct_mention_mode = direct_mention_mode
                orchestrator_instance.direct_mention_target = direct_mention_target
                memory_id = extracted_memory_id
            else:
                orchestrator_instance.session_id = "new_session-12345678901234567890"
                orchestrator_instance.memory_id = "default"
                memory_id = "default"
            
            # RESTORE saved context for the new agent (if available)
            saved_messages = None
            if context_session_key in agent_context_store and agent_name in agent_context_store[context_session_key]:
                saved_messages = agent_context_store[context_session_key][agent_name]
                _flush_log(f"📂 CONTEXT_RESTORE: Found {len(saved_messages)} saved messages for {agent_name}")
            
            # Create agent with restored context
            _flush_log(f"🏗️ AGENT_INVOCATION: Calling create_orchestrator for {agent_name}...")
            try:
                agent = orchestrator_instance.create_orchestrator(
                    orchestrator_instance.session_id, 
                    orchestrator_instance.memory_id, 
                    agent_name,
                    saved_messages=saved_messages
                )
                _flush_log(f"✅ AGENT_INVOCATION: Agent created successfully, type={type(agent)}")
            except Exception as create_err:
                _flush_log(f"❌ AGENT_INVOCATION: create_orchestrator failed: {create_err}", "ERROR")
                import traceback
                _flush_log(f"❌ AGENT_INVOCATION: Traceback: {traceback.format_exc()}", "ERROR")
                raise
            
            current_agent_name = agent_name
            _flush_log(f"✅ Agent created for {agent_name} with session {orchestrator_instance.session_id}")
        else:
            # Agent already exists and type hasn't changed, just update session info if needed
            _flush_log(f"♻️ AGENT_INVOCATION: Reusing existing agent {agent_name}")
            if session_id:
                orchestrator_instance.session_id = session_id
                orchestrator_instance.memory_id = extracted_memory_id
                orchestrator_instance.direct_mention_mode = direct_mention_mode
                orchestrator_instance.direct_mention_target = direct_mention_target
                memory_id = extracted_memory_id
                
                # Update agent state without recreating
                normalized_actor_id = agent_name.replace("_", "-")
                agent.state = {
                    "session_id": session_id,
                    "actor_id": normalized_actor_id,
                    "memory_id": extracted_memory_id,
                }
                _flush_log(f"♻️ Agent reused for {agent_name} with session {session_id}")
        
        if session_id:
            try:
                context_token = set_session_context(session_id)
                _flush_log(f"🔧 AGENT_INVOCATION: Session context set for {session_id}")
            except Exception as ctx_err:
                _flush_log(f"⚠️ AGENT_INVOCATION: Failed to set session context: {ctx_err}", "WARNING")
                context_token = None
        
        # Get tool agent names from config
        stream = payload.get("stream", True)
        _flush_log(f"🎬 AGENT_INVOCATION: stream={stream}, about to invoke agent...")

        if stream:
            _flush_log("🎬 AGENT_INVOCATION: Starting STREAMING mode...")
            try:
                # Build the input for the agent
                agent_input = user_input
                _flush_log(f"🎬 AGENT_INVOCATION: agent_input length={len(str(agent_input))}")

                # If there are file attachments, format them for the agent
                if file_attachments:
                    _flush_log(f"📎 AGENT_INVOCATION: Processing {len(file_attachments)} file attachments...")
                    # Convert to ConverseStream content format
                    content_blocks = []

                    # Process each document attachment
                    document_analyses = []
                    for file_info in file_attachments:
                        try:
                            # Extract S3 location from document block
                            if isinstance(file_info, dict) and "source" in file_info:
                                s3_location = file_info.get("source", {}).get(
                                    "s3Location", {}
                                )
                                s3_uri = s3_location.get("uri", "")

                                if s3_uri.startswith("s3://"):
                                    # Parse S3 URI: s3://bucket/key
                                    s3_parts = s3_uri[5:].split("/", 1)
                                    if len(s3_parts) == 2:
                                        bucket_name = s3_parts[0]
                                        object_key = s3_parts[1]

                                        # Pre-process the document
                                        _flush_log(f"📎 Pre-processing document from s3://{bucket_name}/{object_key}")
                                        analysis = get_s3_as_base64_and_extract_summary_and_facts(
                                            bucket_name, object_key
                                        )

                                        if analysis:
                                            document_name = file_info.get("name", "document")
                                            document_analyses.append(
                                                f"\n\n--- Document: {document_name} ---\n{analysis}"
                                            )
                                            _flush_log(f"📎 Successfully pre-processed document: {document_name}")
                        except Exception as e:
                            _flush_log(f"❌ Failed to pre-process document: {e}", "ERROR")

                    # Append document analyses to user input
                    if document_analyses:
                        enhanced_input = (
                            user_input
                            + "\n\nHere is additional context from attached documents I pre-processed for you:"
                            + "".join(document_analyses)
                        )
                        content_blocks.append({"text": enhanced_input})
                    else:
                        content_blocks.append({"text": user_input})

                    # Add cache point if needed
                    content_blocks.append({"cachePoint": {"type": "default"}})
                    agent_input = content_blocks

                _flush_log(f"🎬 AGENT_INVOCATION: Calling agent.stream_async()...")
                stream_obj = agent.stream_async(agent_input)
                _flush_log(f"🎬 AGENT_INVOCATION: stream_async returned, type={type(stream_obj)}")
                event_count = 0

                events_yielded = False
                _flush_log("🎬 AGENT_INVOCATION: Starting async iteration over stream...")
                async for event in stream_obj:
                    event_count += 1
                    if event_count <= 3:  # Log first few events
                        _flush_log(f"🎬 AGENT_INVOCATION: Event #{event_count}, keys={list(event.keys()) if isinstance(event, dict) else type(event)}")
                    if event.get("message") and event.get("message").get("content"):
                        event["teamName"] = orchestrator_instance.team_name
                        yield event
                        # After stream completes, yield sources as a separate event
                        if collected_sources and collected_sources != {}:
                            _flush_log(f"📦 STREAM: Yielding sources")
                            yield {"type": "sources", "sources": collected_sources}
                
                _flush_log(f"✅ STREAM: Completed with {event_count} events")
                
            except Exception as e:
                _flush_log(f"❌ STREAM: Streaming failed: {e}", "ERROR")
                import traceback
                _flush_log(f"❌ STREAM: Traceback: {traceback.format_exc()}", "ERROR")
                _flush_log("⚠️ STREAM: Falling back to non-streaming mode")
                try:
                    # Build the input for the agent (same as streaming path)
                    agent_input = user_input
                    if file_attachments:
                        content_blocks = [{"text": user_input}]
                        for file_info in file_attachments:
                            content_blocks.append({"document": file_info})
                        content_blocks.append({"cachePoint": {"type": "default"}})
                        agent_input = content_blocks

                    _flush_log(f"🔄 FALLBACK: About to call agent() with input type: {type(agent_input)}")
                    # Fallback to non-streaming response
                    response = agent(agent_input)
                    _flush_log(f"🔄 FALLBACK: agent() returned: {type(response)}")
                    
                    # Safely extract response text
                    try:
                        if hasattr(response, "message") and response.message:
                            content = response.message.get("content")
                            if isinstance(content, list) and len(content) > 0:
                                if isinstance(content[0], dict) and "text" in content[0]:
                                    response_text = content[0]["text"]
                                else:
                                    response_text = str(content[0])
                            elif isinstance(content, str):
                                response_text = content
                            else:
                                response_text = str(content)
                        else:
                            response_text = "No response content available"
                    except (KeyError, IndexError, AttributeError) as e:
                        response_text = f"Error extracting response content: {e}"

                    # Yield the response
                    yield response_text

                    # Yield sources after response
                    if collected_sources:
                        _flush_log(f"📦 FALLBACK: Yielding sources with {len(collected_sources)} sources")
                        yield {"type": "sources", "sources": collected_sources}
                        
                except Exception as fallback_error:
                    _flush_log(f"❌ FALLBACK: Non-streaming failed: {fallback_error}", "ERROR")
                    import traceback
                    _flush_log(f"❌ FALLBACK: Traceback: {traceback.format_exc()}", "ERROR")
                    # Yield an error message
                    yield f"Error processing request: {fallback_error}"
        else:
            # Non-streaming path
            _flush_log("🎬 AGENT_INVOCATION: Starting NON-STREAMING mode...")
            try:
                # Build the input for the agent (same as streaming path)
                agent_input = user_input
                if file_attachments:
                    content_blocks = [{"text": user_input}]
                    for file_info in file_attachments:
                        content_blocks.append({"document": file_info})
                    content_blocks.append({"cachePoint": {"type": "default"}})
                    agent_input = content_blocks

                response = agent(agent_input)
                _flush_log(f"✅ NON-STREAM: Complete response received")

                # Extract the response text
                response_text = ""
                if hasattr(response, "message") and response["message"]:
                    content = response.message.get("content", [])
                    if content and len(content) > 0:
                        response_text = content[0].get("text", "")
                else:
                    response_text = str(response)

                if response_text:
                    # The response may contain agent-message tags from the tools
                    # Split by lines and yield each meaningful chunk
                    lines = response_text.split("\n")
                    current_chunk = ""

                    for line in lines:
                        line_stripped = line.strip()
                        if not line_stripped:
                            continue

                        # Check if this line contains an agent message tag
                        if "<agent-message agent=" in line_stripped:
                            # Yield any accumulated chunk first
                            if current_chunk.strip():
                                yield f"💬 RESPONSE: {current_chunk.strip()}"
                                current_chunk = ""

                            # Yield the agent message as-is
                            yield line_stripped

                        else:
                            # Accumulate regular content
                            current_chunk += line_stripped + " "

                            # Yield when we have a complete thought
                            if (
                                line_stripped.endswith((".", "!", "?"))
                                and len(current_chunk.strip()) > 50
                            ):
                                yield f"💬 RESPONSE: {current_chunk.strip()}"
                                current_chunk = ""

                    # Yield any remaining content
                    if current_chunk.strip():
                        yield f"💬 RESPONSE: {current_chunk.strip()}"

                else:
                    yield "💬 RESPONSE: Analysis completed successfully."

                # Yield sources at the end
                if collected_sources:
                    yield {"type": "sources", "sources": collected_sources}

            except Exception as response_error:
                _flush_log(f"❌ NON-STREAM: Response processing failed: {response_error}", "ERROR")
                import traceback
                _flush_log(f"❌ NON-STREAM: Traceback: {traceback.format_exc()}", "ERROR")
                yield f"❌ ERROR: {response_error}"
            finally:
                # Detach context when done
                try:
                    if context_token:
                        context.detach(context_token)
                        _flush_log(f"🔧 AGENT_INVOCATION: Session context detached")
                except Exception as detach_err:
                    _flush_log(f"⚠️ AGENT_INVOCATION: Failed to detach context: {detach_err}", "WARNING")
    
    except Exception as top_level_error:
        # Top-level exception handler to catch ANY unhandled errors
        _flush_log(f"💥 AGENT_INVOCATION: TOP-LEVEL EXCEPTION: {top_level_error}", "ERROR")
        import traceback
        _flush_log(f"💥 AGENT_INVOCATION: Full traceback:\n{traceback.format_exc()}", "ERROR")
        yield f"❌ FATAL ERROR: {top_level_error}"


# resolver = RuntimeARNResolver(unique_id=os.environ.get("UNIQUE_ID",'1234'), stack_prefix=os.environ.get("STACK_PREFIX",'sim'))
# runtime_url = resolver.resolve_runtime_endpoint(CONFIG.get('agent_id', CONFIG.get('agent_name')))

# host, port = "0.0.0.0", 9000

# Pass runtime_url to http_url parameter AND use serve_at_root=True
# a2a_server = A2AServer(
#    agent=agent,
#    http_url=runtime_url.replace('/invocations',''),
#    serve_at_root=True,  # Serves locally at root (/) regardless of remote URL path complexity
#    port=9000
# )
if __name__ == "__main__":
    app.run()
