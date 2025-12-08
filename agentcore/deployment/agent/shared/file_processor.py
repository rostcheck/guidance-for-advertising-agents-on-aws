"""
File processing utilities for handling S3 documents and converting them for agent analysis.

Supports: PDF, DOCX, PPTX, text files (TXT, MD, CSV, JSON), and images (PNG, JPG, JPEG, WEBP, GIF)
"""

import base64
import logging
import os
import boto3
from typing import List, Optional
from io import BytesIO
from strands import Agent
from strands.models import BedrockModel

logger = logging.getLogger(__name__)


def get_s3_as_base64_and_extract_summary_and_facts(bucket_name: str, object_key: str) -> str:
    """
    Retrieves a document from S3 and extracts content using appropriate method.

    Args:
        bucket_name (str): The name of the S3 bucket.
        object_key (str): The key (path) of the document in the S3 bucket.

    Returns:
        str: Extracted and analyzed content, or error message.
    """
    s3 = boto3.client("s3")
    try:
        # Get the object from S3
        response = s3.get_object(Bucket=bucket_name, Key=object_key)
        file_content = response["Body"].read()

        # Determine file type from extension
        file_ext = object_key.lower().split(".")[-1]
        logger.info(f"Processing {file_ext} file: {object_key}")

        # Route to appropriate processor
        if file_ext == "pdf":
            return process_pdf_document(file_content, object_key)
        elif file_ext == "docx":
            return process_docx_document(file_content, object_key)
        elif file_ext == "pptx":
            return process_pptx_document(file_content, object_key)
        elif file_ext in ["txt", "md", "csv", "json"]:
            return process_text_document(file_content, object_key)
        elif file_ext in ["png", "jpg", "jpeg", "webp", "gif"]:
            return process_image_directly(file_content, object_key, file_ext)
        else:
            return f"Unsupported file type: {file_ext}. Please provide an image or a document of type PDF, DOCX, PPTX, or text files like TXT, MD, CSV, or JSON."

    except Exception as e:
        error_msg = f"Error processing document {object_key}: {str(e)}"
        logger.error(error_msg)
        return error_msg


def process_pdf_document(content: bytes, filename: str) -> str:
    """Process PDF with text extraction and vision fallback"""
    try:
        import PyPDF2

        # Try text extraction first
        pdf_file = BytesIO(content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)

        extracted_text = ""
        for page in pdf_reader.pages:
            extracted_text += page.extract_text() + "\n"

        # Check if we got meaningful text
        if extracted_text.strip() and len(extracted_text.strip()) > 100:
            logger.info(f"PDF text extraction successful: {len(extracted_text)} chars")
            return analyze_extracted_text(extracted_text, "PDF")
        else:
            # Fall back to vision-based processing
            logger.info("PDF text extraction insufficient, using vision method")
            return process_pdf_with_vision(content, filename)

    except Exception as e:
        logger.error(f"PDF processing error: {e}")
        # Final fallback: base64 encoding
        return process_document_as_base64(content, "PDF")


def process_docx_document(content: bytes, filename: str) -> str:
    """Process DOCX by converting to images and using vision AI"""
    try:
        from docx import Document

        # Try text extraction first
        doc = Document(BytesIO(content))
        extracted_text = "\n".join(
            [para.text for para in doc.paragraphs if para.text.strip()]
        )

        # Add table content
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                if row_text.strip():
                    extracted_text += f"\n{row_text}"

        if extracted_text.strip() and len(extracted_text.strip()) > 100:
            logger.info(f"DOCX text extraction successful: {len(extracted_text)} chars")
            return analyze_extracted_text(extracted_text, "Word document")
        else:
            # Fall back to vision processing
            logger.info("DOCX text extraction insufficient, using vision method")
            images = convert_docx_to_images(content)
            return process_images_with_vision(images, filename, "docx")

    except Exception as e:
        logger.error(f"DOCX processing error: {e}")
        return f"Error processing Word document: {str(e)}"


def process_pptx_document(content: bytes, filename: str) -> str:
    """Process PPTX by converting to images and using vision AI"""
    try:
        from pptx import Presentation

        # Try text extraction first
        prs = Presentation(BytesIO(content))
        extracted_text = ""

        for i, slide in enumerate(prs.slides):
            extracted_text += f"\n\n=== Slide {i + 1} ===\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    extracted_text += f"{shape.text}\n"

        if extracted_text.strip() and len(extracted_text.strip()) > 100:
            logger.info(f"PPTX text extraction successful: {len(extracted_text)} chars")
            return analyze_extracted_text(extracted_text, "PowerPoint presentation")
        else:
            # Fall back to vision processing
            logger.info("PPTX text extraction insufficient, using vision method")
            images = convert_pptx_to_images(content)
            return process_images_with_vision(images, filename, "pptx")

    except Exception as e:
        logger.error(f"PPTX processing error: {e}")
        return f"Error processing PowerPoint: {str(e)}"


def process_text_document(content: bytes, filename: str) -> str:
    """Process plain text documents"""
    try:
        text = content.decode("utf-8")
        logger.info(f"Text file processed: {len(text)} chars")
        return analyze_extracted_text(text, "text file")
    except Exception as e:
        logger.error(f"Text processing error: {e}")
        return f"Error processing text file: {str(e)}"


def process_image_directly(content: bytes, filename: str, filetype: str) -> str:
    """Process image files directly with vision AI"""
    try:
        return process_single_image_with_vision(content, 0, filetype)
    except Exception as e:
        logger.error(f"Image processing error: {e}")
        return f"Error processing image file: {str(e)}"


def analyze_extracted_text(text: str, doc_type: str) -> str:
    """Analyze extracted text using Bedrock"""
    try:
        model = BedrockModel(
            model_id="us.anthropic.claude-sonnet-4-20250514-v1:0",
            max_tokens=8000,
            top_p=0.8,
            temperature=0.3,
        )
        analysisAgent = Agent(
            model=model,
            system_prompt=f"You are an expert in analyzing {doc_type} content. Extract key facts, insights, and summarize the main points concisely.",
            tools=[],
            description="Document analysis agent",
        )

        # Limit text length for analysis
        text_to_analyze = text[:15000] if len(text) > 15000 else text
        analysis = analysisAgent(
            f"Analyze this {doc_type} content:\n\n{text_to_analyze}"
        )
        return str(analysis)

    except Exception as e:
        logger.error(f"Text analysis error: {e}")
        return f"Extracted content (analysis failed): {text[:1000]}..."


def process_pdf_with_vision(content: bytes, filename: str) -> str:
    """Process PDF using vision AI for image-based PDFs"""
    try:
        base64_pdf = base64.b64encode(content).decode("utf-8")

        model = BedrockModel(
            model_id="us.anthropic.claude-sonnet-4-20250514-v1:0",
            max_tokens=8000,
            top_p=0.8,
            temperature=0.3,
        )
        analysisAgent = Agent(
            model=model,
            system_prompt="You are an expert in extracting content from document images. Describe all visible text, charts, and visual elements.",
            tools=[],
            description="Vision-based document analyzer",
        )

        analysis = analysisAgent(f"Analyze this PDF document: {base64_pdf}")
        return str(analysis)

    except Exception as e:
        logger.error(f"Vision PDF processing error: {e}")
        return f"Error processing PDF with vision: {str(e)}"


def process_document_as_base64(content: bytes, doc_type: str) -> str:
    """Final fallback: process document as base64"""
    try:
        base64_content = base64.b64encode(content).decode("utf-8")
        return f"Document encoded as base64 ({doc_type}): {base64_content[:200]}... [truncated]"
    except Exception as e:
        return f"Error encoding document: {str(e)}"


def convert_docx_to_images(content: bytes) -> List:
    """Convert DOCX to images for vision processing"""
    try:
        from docx import Document
        import tempfile

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as temp_file:
            temp_file.write(content)
            temp_file.flush()
            temp_path = temp_file.name

        try:
            doc = Document(temp_path)
            images = []

            # Create text-based images from content
            current_content = []
            for para in doc.paragraphs:
                if para.text.strip():
                    current_content.append(para.text)
                    if len(current_content) >= 20:  # Create image every 20 lines
                        img = create_text_image(current_content)
                        images.append(img)
                        current_content = []

            if current_content:
                img = create_text_image(current_content)
                images.append(img)

            return images
        finally:
            if os.path.exists(temp_path):
                os.unlink(temp_path)

    except Exception as e:
        logger.error(f"DOCX to images conversion error: {e}")
        return []


def convert_pptx_to_images(content: bytes) -> List:
    """Convert PPTX to images for vision processing"""
    try:
        from pptx import Presentation
        import tempfile

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
            temp_file.write(content)
            temp_file.flush()
            temp_path = temp_file.name

        try:
            prs = Presentation(temp_path)
            images = []

            for i, slide in enumerate(prs.slides):
                slide_content = [f"=== Slide {i + 1} ==="]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)

                if slide_content:
                    img = create_text_image(slide_content, slide_format=True)
                    images.append(img)

            return images
        finally:
            if os.path.exists(temp_path):
                os.unlink(temp_path)

    except Exception as e:
        logger.error(f"PPTX to images conversion error: {e}")
        return []


def create_text_image(text_lines: List[str], slide_format: bool = False):
    """Create PIL Image from text content"""
    try:
        from PIL import Image, ImageDraw, ImageFont

        width = 1200 if slide_format else 800
        height = max(600, len(text_lines) * 30 + 100)

        img = Image.new("RGB", (width, height), color="white")
        draw = ImageDraw.Draw(img)

        try:
            font = ImageFont.load_default()
        except:
            font = None

        y_offset = 20
        line_height = 22

        for line in text_lines:
            if not line:
                y_offset += line_height // 2
                continue

            # Wrap long lines
            if len(line) > 90:
                words = line.split()
                current_line = ""
                for word in words:
                    test_line = current_line + word + " "
                    if len(test_line) < 90:
                        current_line = test_line
                    else:
                        if current_line.strip():
                            if font:
                                draw.text(
                                    (20, y_offset),
                                    current_line.strip(),
                                    font=font,
                                    fill="black",
                                )
                            else:
                                draw.text(
                                    (20, y_offset), current_line.strip(), fill="black"
                                )
                            y_offset += line_height
                        current_line = word + " "

                if current_line.strip():
                    if font:
                        draw.text(
                            (20, y_offset),
                            current_line.strip(),
                            font=font,
                            fill="black",
                        )
                    else:
                        draw.text((20, y_offset), current_line.strip(), fill="black")
                    y_offset += line_height
            else:
                if font:
                    draw.text((20, y_offset), line, font=font, fill="black")
                else:
                    draw.text((20, y_offset), line, fill="black")
                y_offset += line_height

        return img

    except Exception as e:
        logger.error(f"Text image creation error: {e}")
        # Return minimal fallback image
        from PIL import Image
        return Image.new("RGB", (400, 100), color="white")


def process_images_with_vision(images: List, filename: str, doc_type: str) -> str:
    """Process images using Claude vision"""
    try:
        if not images:
            return f"No images generated from {doc_type} file"

        all_text = []
        for i, image in enumerate(images):
            try:
                page_text = process_single_image_with_vision(image, i + 1, doc_type)
                if page_text.strip():
                    all_text.append(f"=== Page/Slide {i + 1} ===\n{page_text}")
            except Exception as e:
                logger.error(f"Error processing image {i + 1}: {e}")
                all_text.append(f"=== Page/Slide {i + 1} ===\n[Error: {str(e)}]")

        combined_text = "\n\n".join(all_text)
        return analyze_extracted_text(combined_text, doc_type)

    except Exception as e:
        logger.error(f"Vision processing error: {e}")
        return f"Error processing with vision: {str(e)}"


def process_single_image_with_vision(image, page_num: int, doc_type: str) -> str:
    """Process single image with Claude vision"""
    try:
        bedrock_runtime = boto3.client("bedrock-runtime")

        prompt = ''
        response = []
        if doc_type in ['png', 'jpeg', 'jpg', 'webp', 'gif']:
            if doc_type == 'jpg':
                doc_type = 'jpeg'
            prompt = "Analyze this image. If it appears to be a document, extract ALL visible text and describe any visual elements, charts, or diagrams. If it appears to be visual imagery or a creative asset for an ad campaign, return a detailed description of what the image depicts, also noting sentiment, mood, and artistic techniques."
            response = bedrock_runtime.converse(
                modelId="global.anthropic.claude-sonnet-4-5-20250929-v1:0",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"text": prompt},
                            {"image": {"format": doc_type, "source": {"bytes": image}}},
                        ],
                    }
                ],
                inferenceConfig={"maxTokens": 12000, "temperature": 0.8},
            )
        else:
            # Convert PIL Image to bytes
            buffer = BytesIO()
            image.save(buffer, format="PNG", optimize=True, quality=85)
            image_bytes = buffer.getvalue()

            prompt = f"Analyze this {doc_type} page/slide {page_num}. Extract ALL visible text and describe any visual elements, charts, or diagrams."

            response = bedrock_runtime.converse(
                modelId="global.anthropic.claude-sonnet-4-5-20250929-v1:0",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"text": prompt},
                            {"image": {"format": "png", "source": {"bytes": image_bytes}}},
                        ],
                    }
                ],
                inferenceConfig={"maxTokens": 12000, "temperature": 0.1},
            )

        response_text = ""
        if "output" in response and "message" in response["output"]:
            content = response["output"]["message"]["content"]
            for item in content:
                if "text" in item:
                    response_text += item["text"]

        return (
            response_text.strip()
            if response_text
            else f"[No content from page {page_num}]"
        )

    except Exception as e:
        logger.error(f"Single image vision error: {e}")
        return f"[Vision processing failed: {str(e)[:100]}]"
